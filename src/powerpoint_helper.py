from pptx import Presentation 
from pptx.util import Inches 
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from pptx.util import Pt

def addImage(slide, url):
    left = top = Inches(1)
    slide.shapes.add_picture(url, left, top)

""" Ref for slide types: 
0 ->  title and subtitle
1 ->  title and content
2 ->  section header
3 ->  two content
4 ->  Comparison
5 ->  Title only 
6 ->  Blank
7 ->  Content with caption
8 ->  Pic with caption
"""

def add_slide(presentation, type):
    layout = presentation.slide_layouts[type]
    return presentation.slides.add_slide(layout)

def add_splash_slide(presentation, title, subtitle):
    title_slide = add_slide(presentation, 0)
    change_title(title_slide, title, textlayout=Textlayout(68, "Dubai", PP_ALIGN.CENTER))
    title_slide.placeholders[1].text = subtitle


def change_title(slide, title, textlayout):
    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].alignment = textlayout.textalign
    slide.shapes.title.text_frame.paragraphs[0].font.name = textlayout.fontfamily
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(textlayout.fontsize)

def add_text_slide(presentation, title, text):
    slide = add_slide(presentation, 1)
    change_title(slide, title, textlayout=Textlayout(48, "Dubai", PP_ALIGN.LEFT))
    slide.placeholders[1].text = text
    slide.placeholders[1].text_frame.paragraphs[0].font.size = Pt(32)
    return slide

def add_text_slide_with_background(presentation, title, text, background):
    slide = add_slide(presentation, 1)
    addImage(slide, background)
    change_title(slide, title, textlayout=Textlayout(48, "Dubai", PP_ALIGN.LEFT))
    slide.background
    slide.placeholders[1].text = text
    slide.placeholders[1].text_frame.paragraphs[0].font.size = Pt(32)
    return slide

class Textlayout:
    def __init__(self, fontsize, fontfamily, textalign):
        self.fontsize = fontsize
        self.fontfamily = fontfamily
        self.textalign = textalign